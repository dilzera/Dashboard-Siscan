from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(26, 26, 46)
CYAN = RGBColor(23, 162, 184)
PINK = RGBColor(255, 105, 180)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(170, 170, 170)
RED = RGBColor(220, 53, 69)
ORANGE = RGBColor(253, 126, 20)
YELLOW = RGBColor(255, 193, 7)

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PINK
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, bullets=None, highlight=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if bullets:
        y_pos = 1.8
        for bullet in bullets:
            bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(12), Inches(0.6))
            tf = bullet_box.text_frame
            p = tf.paragraphs[0]
            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.font.color.rgb = WHITE
            y_pos += 0.7
    
    if highlight:
        hl_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.5), Inches(12), Inches(1.5))
        hl_box.fill.solid()
        hl_box.fill.fore_color.rgb = RGBColor(40, 40, 70)
        hl_box.line.color.rgb = PINK
        
        hl_text = slide.shapes.add_textbox(Inches(1), Inches(5.7), Inches(11.5), Inches(1.2))
        tf = hl_text.text_frame
        p = tf.paragraphs[0]
        p.text = highlight
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.italic = True
    
    return slide

def add_three_cards_slide(title, cards):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    card_width = 3.8
    gap = 0.4
    start_x = 0.7
    
    for i, card in enumerate(cards):
        x = start_x + i * (card_width + gap)
        
        card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2), Inches(card_width), Inches(4.5))
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = RGBColor(35, 35, 60)
        card_shape.line.color.rgb = RGBColor(60, 60, 90)
        
        card_title = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.3), Inches(card_width - 0.4), Inches(0.8))
        tf = card_title.text_frame
        p = tf.paragraphs[0]
        p.text = card['title']
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.alignment = PP_ALIGN.CENTER
        
        card_desc = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.3), Inches(card_width - 0.4), Inches(2.5))
        tf = card_desc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = card['desc']
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

add_title_slide(
    "Central Inteligente de Câncer de Mama",
    "Curitiba - Paraná\n\nMais que um Dashboard: Uma Plataforma de Inteligência para Jornada do Paciente"
)

add_three_cards_slide("O Problema", [
    {"title": "Tempo de Espera", "desc": "Pacientes aguardam meses entre a solicitação da mamografia e o resultado, perdendo a janela crítica de tratamento precoce."},
    {"title": "Pacientes Perdidas", "desc": "Mulheres com BI-RADS 4 e 5 (alto risco) não retornam para acompanhamento, aumentando mortalidade evitável."},
    {"title": "Dados Fragmentados", "desc": "SISCAN, eSaúde e outros sistemas não conversam entre si, impossibilitando visão completa da jornada da paciente."}
])

add_content_slide("A Solução: Central Inteligente", [
    "Visualização em tempo real de toda a rede de atenção",
    "Rastreamento da jornada completa da paciente",
    "Priorização automática por risco clínico",
    "Alertas inteligentes para busca ativa",
    "Indicadores de desempenho por unidade"
], "Uma plataforma que transforma dados brutos do SISCAN em inteligência acionável.")

add_content_slide("Funcionalidade 1: Navegação da Paciente", [
    "Visualize todo o histórico de uma paciente em uma única tela",
    "Primeiro exame até o mais recente, com evolução do BI-RADS",
    "Tempos de espera entre cada etapa",
    "Unidades de atendimento ao longo da jornada",
    "Filtros por evolução: Positiva (3,4,5 → outros), Negativa (0,1,2 → 3,4,5), Normal"
], "Jornada Completa = Decisões Informadas")

add_three_cards_slide("Funcionalidade 2: Priorização Inteligente", [
    {"title": "CRÍTICO (Vermelho)", "desc": "BI-RADS 5\nAltamente suspeito\n\nSLA: 24 horas"},
    {"title": "ALTO (Laranja)", "desc": "BI-RADS 4\nSuspeito\n\nSLA: 72 horas"},
    {"title": "MÉDIO (Amarelo)", "desc": "BI-RADS 0\nInconclusivo\n\nSLA: 7 dias"}
])

add_content_slide("Funcionalidade 3: Interoperabilidade", [
    "Cruzamento inteligente entre SISCAN e eSaúde",
    "Identificação de pacientes presentes em ambos os sistemas",
    "Detecção de pacientes exclusivas de cada sistema",
    "Alertas de duplicidade por CNS",
    "Garantia de que nenhuma paciente fique 'invisível'"
], "100% de cobertura de dados entre sistemas")

add_content_slide("Funcionalidade 4: Indicadores Clínicos", [
    "COBERTURA: População-alvo atendida vs. meta",
    "AGILIDADE: Tempo médio entre solicitação e resultado",
    "ENCAMINHAMENTOS: Distribuição por categoria BI-RADS",
    "CASOS ESPECIAIS: Fora da faixa etária e situações atípicas"
], "10 indicadores organizados em 4 blocos estratégicos com metas e histórico")

add_content_slide("Funcionalidade 5: Análise por Unidade", [
    "KPIs específicos por unidade de saúde",
    "Heatmap demográfico por faixa etária",
    "Histograma de tempos de espera com meta de 30 dias",
    "Fila de retorno pendente (BI-RADS 0, 3, 4, 5)",
    "Exportação para busca ativa com um clique"
])

add_three_cards_slide("Segurança e Governança", [
    {"title": "Controle de Acesso", "desc": "3 níveis hierárquicos:\nSecretaria, Distrito e Unidade\n\nCada um com visão restrita ao seu escopo"},
    {"title": "Mascaramento de Dados", "desc": "Nomes, CPF, CNS e telefones mascarados por padrão\n\nProteção da privacidade das pacientes"},
    {"title": "Autenticação Segura", "desc": "Senhas criptografadas\nTimeout de sessão\nRecuperação de senha segura"}
])

slide = add_content_slide("Por que é Diferente?", [
    "Dashboard Tradicional: Mostra dados agregados e estáticos",
    "Dashboard Tradicional: Foca em números, não em pessoas",
    "Central Inteligente: Constrói a jornada individual da paciente",
    "Central Inteligente: Prioriza por risco clínico",
    "Central Inteligente: Detecta evolução e alerta automaticamente"
], '"Não contamos exames. Acompanhamos vidas."')

add_content_slide("Visão de Futuro", [
    "INTELIGÊNCIA ARTIFICIAL: Predição de risco de não-retorno",
    "ALERTAS AUTOMATIZADOS: Notificações em tempo real para gestores",
    "APP PARA ACS: Aplicativo móvel para busca ativa em campo",
    "EXPANSÃO NACIONAL: Modelo replicável para outros municípios via DATASUS"
])

add_title_slide(
    "Central Inteligente de Câncer de Mama",
    "Transformando dados em vidas salvas\n\nSecretaria Municipal de Saúde de Curitiba"
)

prs.save('apresentacao_ministerio.pptx')
print("Apresentação salva: apresentacao_ministerio.pptx")
